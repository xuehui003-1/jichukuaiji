#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成《财务机器人课程说课课件》PPTX（16:9）。
v2.0：字号全面放大；受众明确定位为「专家评委 / 同行教师评审」场景；
新增说重难点、说教学过程、说评价、说特色与反思页面。
版本号机制同规划文档：改 VERSION / DATE / REVISION_HISTORY 后重跑，旧版保留。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE

# ============ 每次修订只需要改这里 ============
VERSION = "v2.0"
DATE = "2026-08-26"
REVISION_HISTORY = [
    ("v2.0", "2026-08-26", "字号全面放大；明确定位专家评委/同行评审场景；新增重难点、教学过程、评价、特色与反思页。"),
    ("v1.0", "2026-08-26", "新建：配套《00_财务机器人课程教学规划_v3.0》的说课课件。"),
]
# =============================================

BASE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(BASE)
OUT = os.path.join(ROOT, "01_财务机器人课程说课课件_{}_{}.pptx".format(VERSION, DATE.replace("-", "")))

BLUE = (31, 78, 121)
LIGHT = (217, 226, 243)
GRAY = (89, 89, 89)
DARK = (38, 38, 38)
ACCENT = (198, 89, 17)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def set_run(r, text, size, bold=False, color=DARK, font="微软雅黑"):
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = RGBColor(*color)
    r.font.name = font
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font)


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for line in lines:
        text, size, bold, color = line[0], line[1], line[2], (line[3] if len(line) > 3 else DARK)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(8)
        r = p.add_run()
        set_run(r, text, size, bold, color)
    return tb


def rect(slide, x, y, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(*fill)
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def title_slide(num, title, subtitle=None):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.333, 0.3, BLUE)
    textbox(slide, 0.55, 0.5, 12.2, 1.0, [(title, 34, True, BLUE)])
    rect(slide, 0.55, 1.32, 2.4, 0.06, ACCENT)
    if subtitle:
        textbox(slide, 0.55, 1.55, 12.2, 0.6, [(subtitle, 18, False, GRAY)])
    textbox(slide, 0.55, 7.02, 9, 0.35,
            [("《财务机器人应用与开发》说课 · 文档编号 01 · {} · {}".format(VERSION, DATE), 12, False, GRAY)])
    return slide


def bullets(slide, x, y, w, h, items, size=20, gap=12):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, tuple):
            text, lvl = it
        else:
            text, lvl = it, 0
        prefix = "● " if lvl == 0 else "　– "
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        r = p.add_run()
        set_run(r, prefix + text, size if lvl == 0 else size - 2, lvl == 0)
    return tb


def table(slide, x, y, w, rows, cols, data, col_w=None, font=15, header_fill=BLUE):
    from pptx.util import Emu
    shp = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(0.6 * rows))
    tbl = shp.table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(Inches(w) * cw / total))
    for ri in range(rows):
        for ci in range(cols):
            cell = tbl.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            val = data[ri][ci]
            is_header = ri == 0
            set_run(r, str(val), font - 1 if not is_header else font, is_header,
                    (255, 255, 255) if is_header else DARK)
            if is_header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*header_fill)
    return tbl


# ================= S1 封面 =================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, BLUE)
rect(s, 0, 5.0, 13.333, 0.07, ACCENT)
textbox(s, 1.2, 1.9, 10.9, 2.0, [
    ("财务机器人应用与开发", 48, True, (255, 255, 255)),
    ("说 课 课 件", 44, True, (255, 255, 255)),
], align=PP_ALIGN.CENTER)
textbox(s, 1.2, 4.25, 10.9, 0.7, [
    ("面向专家评委与同行教师 · 教学能力比赛 / 教研评课场景", 20, False, (217, 226, 243)),
], align=PP_ALIGN.CENTER)
textbox(s, 1.2, 5.3, 10.9, 0.9, [
    ("新道智多星 RPA 设计器（离线版） × AI 智能流程自动化", 18, False, (217, 226, 243)),
    ("2026—2027 学年第一学期", 18, False, (217, 226, 243)),
], align=PP_ALIGN.CENTER)
textbox(s, 1.2, 6.6, 10.9, 0.5, [
    ("文档编号 01 ｜ {} ｜ {}".format(VERSION, DATE), 14, False, (180, 195, 215)),
], align=PP_ALIGN.CENTER)

# ================= S2 说课框架 =================
s = title_slide(1, "说课框架", "本课件定位：专家评委 / 同行教师评审场景")
bullets(s, 0.9, 2.2, 11.5, 4.5, [
    "说教材 —— 课程定位与依据",
    "说学情 —— 已经会什么、缺什么",
    "说目标 —— 知识 · 技能 · 素养",
    "说重难点 —— 重点与突破策略",
    "说内容 —— 原有六项目 + 新增项目七",
    "说教法学法 —— 任务驱动 + AI 点睛",
    "说教学过程 —— 90 分钟课堂标准流程",
    "说评价与反思 —— 过程性 + 作品化",
], size=21, gap=12)

# ================= S3 说教材 =================
s = title_slide(2, "说教材：课程定位与依据", "培养能指挥「数字员工」的财务人才")
bullets(s, 0.9, 2.2, 11.5, 4.5, [
    "课程性质：会计类专业课程，衔接《基础会计》凭证、账簿、报表知识",
    "课程依据：会计专业人才培养方案 · 数字化财务岗位能力要求",
    "课程理念：RPA 是手，AI 是脑 —— 机器人干活，AI 想事，学生指挥",
    "工具体系：智多星 RPA 设计器（离线版）+ 云端 AI 平台 + 飞书多维表格",
], size=21, gap=14)

# ================= S4 说学情 =================
s = title_slide(3, "说学情：已经会什么、缺什么", "零基础学生的破局之道")
bullets(s, 0.9, 2.2, 11.5, 4.5, [
    "已学《基础会计》：看得懂凭证、账簿、报表 —— 财务语言无障碍",
    "工具习惯：熟悉企业微信、飞书；亲近短视频与 AI 生成内容",
    "AI 与编程零基础：好奇与畏难并存",
    "对策：先看热闹、再学门道 —— 每课 10–15 分钟「AI 点睛」降门槛",
], size=21, gap=14)

# ================= S5 说目标 =================
s = title_slide(4, "说教学目标", "三维目标 · 可检测")
table(s, 0.9, 2.3, 11.5, 4, 2, [
    ("维度", "具体目标"),
    ("知识", "理解 RPA 与 AI 在财务中的作用；掌握智多星变量、命令与自动化组件"),
    ("技能", "能开发财务机器人；能搭建财务问答智能体、建设知识库、用 AI 整理财务资料"),
    ("素养", "数据安全（凭证脱敏）；人机协同；复核与批判思维 —— 机器人和 AI 都可能错"),
], col_w=[2.2, 9.8], font=16)

# ================= S6 说重难点 =================
s = title_slide(5, "说重难点与突破策略", "重点抓得准，难点破得开")
bullets(s, 0.9, 2.2, 11.5, 4.5, [
    "重点：变量与命令、Excel / 浏览器 / 邮件自动化、发票查验机器人",
    "难点：循环与条件逻辑、异常处理、AI 幻觉识别与人工复核",
    "突破策略：任务小步拆解 → 报错现场教学 → 人机 PK 与双人复核游戏",
], size=21, gap=14)

# ================= S7 说内容（上） =================
s = title_slide(6, "说内容（上）：原有六大项目 · 26 课时", "沿用智多星设计器主线，内容与顺序不变")
table(s, 0.9, 2.2, 11.5, 7, 3, [
    ("项目", "主要内容", "课时"),
    ("项目一 说课与 RPA 认知", "说课；AI 发展与职业规划；RPA 介绍与商科应用", "3"),
    ("项目二 初识 RPA", "基本介绍；工具介绍；流程优化；注册安装", "4"),
    ("项目三 变量与命令", "五类变量；命令含义、区分与应用场景", "6"),
    ("项目四 自动化处理", "Excel / 浏览器 / 邮件自动化；屏幕录制；OCR", "6"),
    ("项目五 实施方法论", "需求评估 → 方案设计 → 配置开发 → 改进", "2"),
    ("项目六 税务机器人", "验证码识别；发票查验；企业信息查验；发票开具", "5"),
], col_w=[3.4, 6.8, 1.2], font=15)

# ================= S8 说内容（下） =================
s = title_slide(7, "说内容（下）：新增项目七 · 6 课时", "今年新增：让「RPA 的手」接上「AI 的脑」")
table(s, 0.9, 2.2, 11.5, 5, 3, [
    ("任务", "内容", "课时"),
    ("任务 1", "生成式 AI 与财务提示词：原理、四要素、财务场景练习、伦理与数据安全", "2"),
    ("任务 2", "财务知识库建设：报销制度 / 税法要点 → 上传切片 → 问答验证", "1"),
    ("任务 3", "智能体财务助手：搭「财务政策问答助手」，演示智多星取数 + 助手解读", "2"),
    ("任务 4", "AIGC 整理与多维表格：单据清单、分析报告、飞书多维表格 AI 字段", "1"),
], col_w=[1.6, 9.0, 1.2], font=15)

# ================= S9 说教法学法 =================
s = title_slide(8, "说教法与学法", "四原则：扣题 · 有来路 · 有逻辑链 · 可复现")
bullets(s, 0.9, 2.2, 11.5, 4.5, [
    "教法：任务驱动 + 教师示范 + 费曼四步（现象 → 猜测 → 讲解 → 练习）",
    "学法：做中学、小组协作、作品展示互评",
    "每课「AI 点睛」10–15 分钟，全部紧扣财务机器人主题",
    "创意任务库：配音连环画、新品发布会、抓 bug 大赛、流水线接力赛",
], size=21, gap=14)

# ================= S10 说教学过程 =================
s = title_slide(9, "说教学过程：90 分钟标准课堂", "每个环节都有明确的教学意图")
table(s, 0.9, 2.2, 11.5, 8, 3, [
    ("时段", "环节", "教学意图"),
    ("0–5 min", "引入：财务场景问题（漫画 / 视频）", "立靶子：为什么要机器人"),
    ("5–20 min", "演示：教师当堂搭流程", "先见完整成品"),
    ("20–23 min", "暂停点：先猜后讲", "预测效应，记得牢"),
    ("23–38 min", "讲解：核心概念拆解", "循环 / 条件等关键概念"),
    ("38–55 min", "任务 1：小组协作改造流程", "从模仿到改造"),
    ("55–80 min", "任务 2 + 展示：人机 PK、互评投票", "展示即考核"),
    ("80–90 min", "收尾：坑点小结 + 课后任务", "防课后翻车"),
], col_w=[2.0, 5.4, 4.2], font=14)

# ================= S11 说评价 =================
s = title_slide(10, "说评价：过程性 + 作品化")
table(s, 0.9, 2.3, 11.5, 4, 3, [
    ("类别", "考核项目", "占比"),
    ("过程性", "课堂任务积分 + 作品入班级彩蛋库", "40%"),
    ("过程性", "机器人作业：需求说明书、抓 bug 诊断书、机器人简历", "20%"),
    ("期末", "新品发布会：机器人演示 + AI 助手演示 + 答辩", "40%"),
], col_w=[2.2, 8.0, 1.6], font=16)

# ================= S12 说特色与反思 =================
s = title_slide(11, "说特色与教学反思")
bullets(s, 0.9, 2.2, 11.5, 4.5, [
    "特色一 主线稳：智多星六大项目全覆盖，与原有课程无缝衔接",
    "特色二 点睛活：每课 AI 点睛紧扣财务机器人，逻辑链完整",
    "特色三 作品即作业：配音连环画、发布会、抓 bug 大赛",
    "反思预案一：机房无外网 → 教师演示 + 学生手机课后完成",
    "反思预案二：AI 结果不可控 → 双人复核机制贯穿全程",
], size=21, gap=14)

# ================= S13 结束页 =================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, BLUE)
textbox(s, 1.2, 2.9, 10.9, 1.4, [
    ("敬请各位领导、专家、老师指导", 38, True, (255, 255, 255)),
], align=PP_ALIGN.CENTER)
textbox(s, 1.2, 4.6, 10.9, 0.5, [
    ("让机器人干活，让 AI 想事，让学生爱上指挥", 20, False, (217, 226, 243)),
], align=PP_ALIGN.CENTER)

prs.save(OUT)
print("已生成:", OUT)
